import os
import fitz
import re
import docx2txt
import openpyxl
from pptx import Presentation
from odf import text, teletype
from odf.opendocument import load
from bs4 import BeautifulSoup
import requests
import json
# Solr URL
class SolrProcessor:
    def __init__(self, solr_url):
        self.solr_url = solr_url
        self.check()
    def check(self) :
        data  = self.file_to_jsondata()

        headers = {"Content-Type": "application/json"}
        # 傳送 POST 請求到 Solr
        response = requests.post( self.solr_url+'/select?q=*:*&rows=0', headers=headers)
        res = response.json()
        if(res['response']['numFound'] != len(data)) : 
            print(res['response']['numFound'] ,len(data) ) 
            self.delete_solr_data()         
            self.data_to_solr()
    def read_docx(self,file_path):
        content = docx2txt.process(file_path )
        content = re.sub(r"\s+", "", content)
        return content
    def read_pdf(self,pdf_file):
        text_content = ""
        doc = fitz.open(pdf_file)
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num) 
            text_content += page.get_text("text")
            
        doc.close()
        text_content = re.sub(r"\s+", "", text_content)
        return text_content
    def read_pptx(self,ppt_file):
        presentation = Presentation(ppt_file)

        content = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text
        content = re.sub(r"\s+", "", content)
        return content
    def read_xlsx(xlxs_file):
        workbook = openpyxl.load_workbook(xlxs_file)
        content = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    if(cell.value ):   
                        content+=str(cell.value)
        return content
    def read_odt(self,odt_file):
        doc = load(odt_file)
        all_paragraphs = []
        for body in doc.getElementsByType(text.P):
            paragraph_text = teletype.extractText(body)
            all_paragraphs.append(paragraph_text)
        content = "".join(all_paragraphs)
        content = re.sub(r"\s+", "", content)
        return content
    def read_html(self,html_file) :
        with open(html_file, "r", encoding="utf-8") as file:
            html_content = file.read()
            soup = BeautifulSoup(html_content, "html.parser")
        body = soup.find("body")
        all_text = ''
        if(body):
            all_text = body.get_text(separator=' ')
            all_text = ' '.join(all_text.split())  # 移除多餘的空白字符
        return all_text
    def file_to_jsondata(self):

        current_directory = os.getcwd()  # 取得目前的資料夾路徑
        data = []
        for root, dirs, files in os.walk(current_directory):
            if os.path.basename(root) == "myenv":
                continue
            for file in files:
                filename, file_extension = os.path.splitext(file)
                file_extension = file_extension[1:]  # 去除副檔名前面的點
                file_path = os.path.join(root, file)
                file_data = {}
                file_data["File_Name_txt_ja"] = filename
                file_data["File_Path_txt_ja"] = file_path
                file_data["File_Content_txt_ja"] = ""
                content = ""      
                if(file_extension == 'pdf') :
                    content = self.read_pdf(file_path)
                elif(file_extension == 'odt') :
                    content = self.read_odt(file_path)             
                elif(file_extension == 'html') :
                    content = self.read_html(file_path)
                elif(file_extension == 'docx') :
                    content = self.read_docx(file_path)     
                elif(file_extension == 'pptx') :
                    content = self.read_pptx(file_path)    
                elif(file_extension == 'xlxs') :
                    content = self.read_xlsx(file_path)
                if content :
                    file_data["File_Content_txt_ja"] = content
                data.append(file_data)               
        return data
    def delete_solr_data(self):
        # 設定刪除所有文件的 HTTP 請求
        delete_all_request_data = {
            "delete": {"query": "*:*"}
        }

        # 發送刪除請求
        response = requests.post(self.solr_url+'/update?commit=true', json=delete_all_request_data)

        # 檢查回應狀態碼
        if response.status_code == 200:
            print("所有文件已從 Solr 刪除")
        else:
            print("刪除失敗，回應狀態碼：", response.status_code)    
    def data_to_solr(self) :
        json_data = self.file_to_jsondata()
        # 將 JSON 資料轉換為字串
        json_string = json.dumps(json_data)

        # 設定請求標頭
        headers = {"Content-Type": "application/json"}

        # 傳送 POST 請求到 Solr
        response = requests.post(self.solr_url+"/update?commit=true", data=json_string, headers=headers)

        # 檢查回應狀態碼
        if response.status_code == 200:
            print("上傳成功！")
        else:
            print("上傳失敗：", response.text)
    def search_solr(self,query):
        params = {
            'q': 'File_Content_txt_ja : "'+query +'"',
            'wt': 'json',
            "rows":"62400",
            "hl.tag.pre":"<span class=\"text-danger\">",
            "hl":"true",
            "indent":"true",
            "hl.fragsize":"200",
            "q.op":"OR",
            'hl.bs.type': 'WORD',
            'hl.fragAlignRatio': "0.5",
            "hl.fl":"File_Content_txt_ja",
            "hl.method":"unified",
            "hl.tag.post":"</span>"
        }

        response = requests.get(self.solr_url+'/select', params=params)

        if response.status_code == 200:
            return response.json()
        else:
            print(f"Error: {response.text}")
            return None